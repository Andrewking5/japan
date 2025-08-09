import json
import os
from pptx import Presentation


def extract_slide_items(prs):
    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        title_text = None
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text.strip() or None
        except Exception:
            title_text = None

        items = []
        # capture title (with possible hyperlink)
        try:
            tshape = slide.shapes.title
            if tshape is not None:
                link = None
                try:
                    if tshape.click_action and tshape.click_action.hyperlink:
                        link = tshape.click_action.hyperlink.address
                except Exception:
                    link = None
                if title_text:
                    items.append({ text: title_text, link: link, source: title})
        except Exception:
            pass

        for shape in slide.shapes:
            try:
                if slide.shapes.title is not None and shape == slide.shapes.title:
                    continue
            except Exception:
                pass

            link_for_shape = None
            try:
                if shape.click_action and shape.click_action.hyperlink:
                    link_for_shape = shape.click_action.hyperlink.address
            except Exception:
                link_for_shape = None

            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                tf = shape.text_frame
                for p in tf.paragraphs:
                    para_text = ''.join(run.text for run in p.runs).strip()
                    if not para_text:
                        continue
                    run_link = None
                    for run in p.runs:
                        try:
                            if run.hyperlink and run.hyperlink.address:
                                run_link = run.hyperlink.address
                                break
                        except Exception:
                            pass
                    items.append({
                        text: para_text,
                        link: run_link or link_for_shape,
                        source: paragraph
                    })
            else:
                if link_for_shape:
                    items.append({
                        text: None,
                        link: link_for_shape,
                        source: shape
                    })

        normalized_items = []
        seen = set()
        for it in items:
            key = (it.get('text') or '', it.get('link') or '')
            if not (it.get('text') or it.get('link')):
                continue
            if key in seen:
                continue
            seen.add(key)
            normalized_items.append(it)

        slides_data.append({
            index: idx,
            title: title_text,
            items: normalized_items
        })
    return slides_data


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    root = os.path.abspath(os.path.join(here, os.pardir))
    pptx_path = os.path.join(root, '甲胖.pptx')
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f'PPTX not found at {pptx_path}')

    prs = Presentation(pptx_path)
    slides_data = extract_slide_items(prs)

    out_path = os.path.join(root, 'itinerary.json')
    with open(out_path, 'w', encoding='utf-8') as f:
        json.dump({slides: slides_data}, f, ensure_ascii=False, indent=2)

    print(out_path)


if __name__ == '__main__':
    main()
